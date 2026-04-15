import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Attempt to handle a known python-pptx bug with newer python versions
    try:
        collections.Sequence = collections.abc.Sequence
    except AttributeError:
        pass
        
    prs = Presentation()

    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Haven Chat"
    subtitle.text = "Production-Grade AI Moderated Chat Architecture\nTechnical Analysis & Overview"

    # --- Slide 2: Project Overview ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Project Overview"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Haven is designed to be a sanctuary from toxic internet interactions."
    p = tf.add_paragraph()
    p.text = "Ephemeral Real-Time Messaging: Messages appear and expire dynamically."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "AI-Powered Moderation: Instantaneous filtering of harmful content."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Privacy First: Secure authentication and row-level security."
    p.level = 1

    # --- Slide 3: Tech Stack ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Core Technology Stack"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "The platform is separated into distinct micro-services:"
    
    p = tf.add_paragraph()
    p.text = "Frontend: Next.js 16 (App Router), React 19, Tailwind CSS, shadcn/ui"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Database & Auth: Supabase (PostgreSQL, Realtime WebSockets)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "ML Service: FastAPI (Python), Scikit-Learn, Google Gemini API"
    p.level = 1

    # --- Slide 4: Hybrid Moderation Engine ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "The Hybrid AI Pipeline"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "A multi-stage decision-fusion engine optimizes cost and speed."

    p = tf.add_paragraph()
    p.text = "Stage 1 (Local): TF-IDF + Logistic Regression evaluates input natively (<50ms)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Stage 2 (Fallback): If probability is severely borderline (0.45 - 0.85), Gemini evaluates context."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Stage 3 (Bypass): Whitelists instantly approve gaming terms (e.g. 'kill the process')."
    p.level = 1

    # --- Slide 5: False-Positive Mitigation ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "False-Positive Reduction Strategy"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Ensuring safe messages aren't accidentally trapped."

    p = tf.add_paragraph()
    p.text = "Multilingual Awareness: Model handles Hinglish slurs ('chutiya', 'pagal') securely."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Context Injection: Gemini analyzes the last 5 chat messages, not just isolated sentences."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Safety Net: If Gemini's API fails, it structurally defaults to a safe hard-block."
    p.level = 1

    # --- Slide 6: Database & Flow ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Data Flow & Security"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Modern infrastructure handles the asynchronous data."

    p = tf.add_paragraph()
    p.text = "Clients post strictly to the ML endpoint before database insertion."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Supabase Row-Level-Security (RLS) policies prevent unauthorized reads/writes."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Optimistic UI rendering mimics FAANG-level latency for end-users."
    p.level = 1

    prs.save('Haven_Architecture_Presentation.pptx')
    print("PowerPoint presentation 'Haven_Architecture_Presentation.pptx' created successfully.")

if __name__ == '__main__':
    create_presentation()
